"""证书批量生成器 — 网页版"""

import io
import os
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path

import pandas as pd
import streamlit as st
from pptx import Presentation
from pypinyin import lazy_pinyin, load_phrases_dict


# ====================== 初始化 ======================

def _init_surname_dict():
    csv_path = Path(__file__).parent / "surname_pinyin.csv"
    if not csv_path.exists():
        return
    try:
        df = pd.read_csv(str(csv_path), dtype=str)
        phrases = {}
        for _, row in df.iterrows():
            char = str(row.get("character", "")).strip()
            pinyin_str = str(row.get("surname_pinyin", "")).strip()
            if char and pinyin_str:
                phrases[char] = [[s] for s in pinyin_str.split()]
        if phrases:
            load_phrases_dict(phrases)
    except Exception:
        pass


def _check_deps():
    return shutil.which("soffice") is not None and shutil.which("pdftoppm") is not None


_init_surname_dict()
DEPS_OK = _check_deps()

BATCH_SIZE = 15


# ====================== 工具函数 ======================

def is_chinese(text):
    if not text:
        return False
    return any('一' <= ch <= '鿿' for ch in str(text))


def chinese_to_pinyin(name):
    if not name or str(name).strip() == "":
        return ""
    return " ".join(w.capitalize() for w in lazy_pinyin(str(name), strict=False))


def clean_cell(val):
    s = str(val).strip()
    return s[:-2] if s.endswith(".0") else s


def fill_template(prs, replacements):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for key, value in replacements.items():
                            run.text = run.text.replace(key, value)
    return prs


def pptx_bytes_to_png(pptx_bytes):
    with tempfile.TemporaryDirectory() as tmpdir:
        base = Path(tmpdir)
        pptx_path = base / "input.pptx"
        pptx_path.write_bytes(pptx_bytes)

        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             "--outdir", str(base), str(pptx_path)],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            return None

        pdf_files = list(base.glob("*.pdf"))
        if not pdf_files:
            return None

        prefix = base / "page"
        result = subprocess.run(
            ["pdftoppm", "-png", "-r", "200", "-f", "1", "-l", "1",
             str(pdf_files[0]), str(prefix)],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode != 0:
            return None

        generated = list(base.glob("page-*.png"))
        return generated[0].read_bytes() if generated else None


# ====================== 页面 ======================

st.set_page_config(page_title="证书批量生成器", page_icon="📜", layout="wide")

if not DEPS_OK:
    st.error("❌ 服务端未安装 LibreOffice / poppler，无法使用。")
    st.stop()

# 初始化 session state
for key, default in [
    ("_trigger", False),
    ("_total_rows", []),
    ("_pptx_bytes", None),
    ("_has_en", False),
    ("_batch_idx", 0),
    ("_total", 0),
    ("_work_dir", None),
    ("_png_names", []),
    ("_first_name", None),
    ("_ok", 0),
    ("_fail", 0),
]:
    if key not in st.session_state:
        st.session_state[key] = default

# ===== 标题 =====
st.title("📜 证书批量生成器")
st.caption("上传 CSV + PPTX 模板，批量导出证书 PNG 图片")
st.divider()

left, right = st.columns([1, 1])

# ===== 左栏：配置 =====
with left:
    with st.expander("📖 使用说明", expanded=True):
        st.markdown("""
        **CSV 文件格式**（必须有表头）：
        | number | id | cn_name | en_name |
        |--------|----|---------|---------|
        | 001 | 2026001 | 张三 | |
        | 002 | 2026002 | Lisa | |
        | 003 | 2026003 | 卜苗 | |

        **PPTX 模板占位符**：`{{ID}}` `{{CN_NAME}}` `{{EN_NAME}}`

        **EN_NAME 规则**：
        - cn_name 是英文 → 留空
        - en_name 列有值 → 直接用
        - 否则 → 拼音自动生成（含多音字纠正）
        """)

    ul1, ul2 = st.columns(2)
    with ul1:
        csv_file = st.file_uploader("📄 上传 CSV", type=["csv"], key="csv_up")
    with ul2:
        pptx_file = st.file_uploader("📄 上传 PPTX", type=["pptx"], key="pptx_up")

    if csv_file is not None:
        try:
            csv_bytes_preview = csv_file.getvalue()
            preview_df = pd.read_csv(io.BytesIO(csv_bytes_preview),
                                     dtype={"id": str, "number": str})
            st.caption(f"📋 数据预览（共 {len(preview_df)} 条，显示前 3 条）")
            st.dataframe(preview_df.head(3), use_container_width=True)
        except Exception as e:
            st.error(f"CSV 解析失败: {e}")

    if csv_file and pptx_file:
        if st.button("🚀 开始批量生成 PNG", type="primary", use_container_width=True):
            # 准备批次数据
            csv_bytes = csv_file.getvalue()
            pptx_bytes = pptx_file.getvalue()
            df = pd.read_csv(io.BytesIO(csv_bytes), dtype={"id": str, "number": str})

            # 清理旧目录
            if st.session_state._work_dir:
                shutil.rmtree(st.session_state._work_dir, ignore_errors=True)

            st.session_state._pptx_bytes = pptx_bytes
            st.session_state._total_rows = df.to_dict("records")
            st.session_state._total = len(df)
            st.session_state._has_en = "en_name" in df.columns
            st.session_state._batch_idx = 0
            st.session_state._png_names = []
            st.session_state._first_name = None
            st.session_state._ok = 0
            st.session_state._fail = 0
            st.session_state._work_dir = str(Path(tempfile.mkdtemp(prefix="cert_")))
            st.session_state._trigger = True
            st.rerun()

# ===== 右栏：进度 / 下载 / 预览 =====
with right:
    if st.session_state._trigger:
        _total = st.session_state._total
        _batch_idx = st.session_state._batch_idx
        _total_rows = st.session_state._total_rows
        _pptx_bytes = st.session_state._pptx_bytes
        _has_en = st.session_state._has_en
        _work_dir = Path(st.session_state._work_dir)

        batch_end = min(_batch_idx + BATCH_SIZE, _total)
        is_last = (batch_end >= _total)
        batch_rows = _total_rows[_batch_idx:batch_end]

        progress_bar = st.progress(
            _batch_idx / _total if _total > 0 else 0,
            text=f"第 {_batch_idx + 1}-{batch_end} 张（共 {_total}）"
        )
        status_text = st.empty()
        preview_spot = st.empty()

        # 第一批时展示预估时间
        if _batch_idx == 0 and _total > BATCH_SIZE:
            est_min = _total * 3 // 60
            st.caption(f"预计 {est_min} 分钟左右，每批 {BATCH_SIZE} 张自动接力")

        try:
            for j, row in enumerate(batch_rows):
                uid = clean_cell(row.get("id", ""))
                cn_name = str(row.get("cn_name", "")).strip()
                num = clean_cell(row.get("number", ""))

                en_name = ""
                if is_chinese(cn_name):
                    if _has_en:
                        val = row.get("en_name", "")
                        if pd.notna(val) and str(val).strip() != "":
                            en_name = str(val).strip()
                    if en_name == "":
                        en_name = chinese_to_pinyin(cn_name)

                prs = Presentation(io.BytesIO(_pptx_bytes))
                fill_template(prs, {
                    "{{ID}}": uid,
                    "{{CN_NAME}}": cn_name,
                    "{{EN_NAME}}": en_name
                })
                buf = io.BytesIO()
                prs.save(buf)

                png_bytes = pptx_bytes_to_png(buf.getvalue())
                if png_bytes:
                    fname = f"{num}_{uid}_{cn_name}.png"
                    (_work_dir / fname).write_bytes(png_bytes)
                    st.session_state._png_names.append(fname)
                    st.session_state._ok += 1

                    # 第一张立即展示预览
                    if st.session_state._first_name is None:
                        st.session_state._first_name = fname
                        preview_spot.image(
                            str(_work_dir / fname),
                            caption=fname,
                            use_container_width=True
                        )
                else:
                    st.session_state._fail += 1

                global_i = _batch_idx + j + 1
                progress_bar.progress(
                    global_i / _total,
                    text=f"正在处理 {global_i}/{_total}"
                )
                status_text.caption(f"当前：{cn_name}")

        except Exception as e:
            st.session_state._trigger = False
            st.error(f"生成失败: {e}")
            st.rerun()

        progress_bar.empty()
        status_text.empty()

        if is_last:
            preview_spot.empty()
            st.session_state._trigger = False
        else:
            st.session_state._batch_idx = batch_end
        st.rerun()

    elif st.session_state._work_dir:
        # 生成完成，展示下载和预览
        _work_dir = Path(st.session_state._work_dir)
        ok = st.session_state._ok
        fail = st.session_state._fail
        png_names = st.session_state._png_names
        first_name = st.session_state._first_name

        if ok > 0:
            loading = st.empty()
            loading.info("⏳ 正在打包 ZIP 文件，请稍候...")
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
                for name in png_names:
                    zf.write(_work_dir / name, arcname=name)
            zip_buf.seek(0)
            loading.empty()

            st.download_button(
                label=f"⬇ 下载全部（{ok} 张，ZIP）",
                data=zip_buf.getvalue(),
                file_name="certificates.zip",
                mime="application/zip",
                use_container_width=True
            )
            st.caption("如果下载未弹出，请允许浏览器下载多个文件")

            msg = f"✅ 成功 {ok} 张"
            if fail > 0:
                msg += f"，失败 {fail} 张"
            st.success(msg)

            if first_name:
                st.image(
                    str(_work_dir / first_name),
                    caption=first_name,
                    use_container_width=True
                )
        else:
            st.error(f"❌ 全部 {fail} 张转换失败")

    else:
        st.markdown("""
        <div style="height:300px;display:flex;align-items:center;justify-content:center;
        color:#aaa;font-size:16px;text-align:center;border:2px dashed #ddd;border-radius:12px">
        👈 上传 CSV 和 PPTX 模板<br>点击生成后，证书预览将显示在这里
        </div>
        """, unsafe_allow_html=True)

st.divider()
st.caption("💡 姓氏多音字自动纠正（卜→Bǔ、单→Shàn、仇→Qiú 等）| 基于 pypinyin")
