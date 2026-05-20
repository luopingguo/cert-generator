#!/usr/bin/env python3
"""
证书批量生成器
=======================
根据 CSV 数据 + PPTX 模板，批量生成带参数的 PPTX 文件，并可转换为 PNG。

用法:
    python run.py                          # 使用默认配置
    python run.py --csv mydata.csv         # 指定 CSV
    python run.py --no-png                 # 只生成 PPTX，不转 PNG

依赖:
    pip install pandas python-pptx pypinyin tqdm
    (可选) macOS + LibreOffice + poppler 用于 PPTX → PNG 转换

CSV 格式:
    number,id,cn_name,en_name
    001,2026001,张三,
    002,2026002,Lisa,
    003,2026003,卜苗,

PPTX 模板占位符:
    {{ID}}       → 替换为 id 列
    {{CN_NAME}}  → 替换为 cn_name 列
    {{EN_NAME}}  → 替换逻辑见下方说明

EN_NAME 填充规则:
    1. cn_name 本身是英文 → 留空（不填）
    2. CSV 中 en_name 有内容 → 直接用
    3. 否则 → 拼音自动生成（含姓氏多音字纠正）
"""

import argparse
import os
import subprocess
import sys
import tempfile
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pypinyin import lazy_pinyin, load_phrases_dict

# ====================== 配置（按需修改） ======================
TEMPLATE_PPT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "template.pptx")
CSV_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "sample.csv")
SURNAME_CSV = os.path.join(os.path.dirname(os.path.abspath(__file__)), "surname_pinyin.csv")
OUTPUT_PPT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output/pptx")
OUTPUT_PNG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output/png")

# LibreOffice 路径（仅 macOS PNG 转换时需要）
SOFFICE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
POPPLER_BIN = Path("/opt/homebrew/opt/poppler/bin")
# =============================================================


def is_chinese(text):
    """判断字符串是否包含中文字符"""
    if not text:
        return False
    return any('一' <= ch <= '鿿' for ch in str(text))


def load_surname_dict(csv_path):
    """从 CSV 加载姓氏多音字字典"""
    if not os.path.exists(csv_path):
        print(f"[警告] 未找到姓氏字典 {csv_path}，跳过多音字纠正")
        return {}
    df = pd.read_csv(csv_path, dtype=str)
    phrases = {}
    for _, row in df.iterrows():
        char = row["character"].strip()
        pinyin_str = row["surname_pinyin"].strip()
        phrases[char] = [[s] for s in pinyin_str.split()]
    return phrases


def chinese_to_pinyin(name):
    """中文名 → 拼音（首字母大写，空格分隔）"""
    if not name or str(name).strip() == "":
        return ""
    pinyin_list = lazy_pinyin(str(name), strict=False)
    return " ".join(w.capitalize() for w in pinyin_list)


def fill_template(pptx_path, replacements):
    """替换 PPTX 中的占位符，保留原格式"""
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for key, value in replacements.items():
                            run.text = run.text.replace(key, value)
    return prs


def convert_to_png(pptx_path, output_png, pdftoppm):
    """单页 PPTX → PNG（依赖 LibreOffice + poppler）"""
    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_dir = Path(tmpdir)

        result = subprocess.run(
            [SOFFICE_PATH, "--headless", "--convert-to", "pdf",
             "--outdir", str(pdf_dir), str(pptx_path)],
            capture_output=True, text=True
        )
        if result.returncode != 0:
            return False

        pdf_files = list(pdf_dir.glob("*.pdf"))
        if not pdf_files:
            return False

        prefix = Path(tmpdir) / "page"
        result = subprocess.run(
            [str(pdftoppm), "-png", "-r", "200", "-f", "1", "-l", "1",
             str(pdf_files[0]), str(prefix)],
            capture_output=True, text=True
        )
        if result.returncode != 0:
            return False

        generated = list(Path(tmpdir).glob("page-*.png"))
        if not generated:
            return False

        generated[0].rename(output_png)
        return True


def resolve_en_name(cn_name, en_name_raw, has_en_column):
    """按规则决定 EN_NAME 的值"""
    if not is_chinese(cn_name):
        return ""                               # 英文名 → 不填
    if has_en_column and pd.notna(en_name_raw) and str(en_name_raw).strip() != "":
        return str(en_name_raw).strip()         # CSV 有值 → 直接用
    return chinese_to_pinyin(cn_name)           # 拼音自动生成


def clean_cell(val):
    """去除 pandas 读数的 .0 后缀"""
    s = str(val).strip()
    return s[:-2] if s.endswith(".0") else s


def main():
    parser = argparse.ArgumentParser(description="证书批量生成器")
    parser.add_argument("--csv", default=CSV_FILE, help="CSV 数据文件路径")
    parser.add_argument("--template", default=TEMPLATE_PPT, help="PPTX 模板路径")
    parser.add_argument("--no-png", action="store_true", help="跳过 PNG 转换")
    parser.add_argument("--out-pptx", default=OUTPUT_PPT_DIR, help="PPTX 输出目录")
    parser.add_argument("--out-png", default=OUTPUT_PNG_DIR, help="PNG 输出目录")
    args = parser.parse_args()

    # 加载姓氏字典
    load_phrases_dict(load_surname_dict(SURNAME_CSV))

    # 检查必需文件
    for f, label in [(args.template, "模板"), (args.csv, "CSV")]:
        if not os.path.exists(f):
            sys.exit(f"[错误] 找不到{label}文件: {f}")

    os.makedirs(args.out_pptx, exist_ok=True)

    # ========== 阶段 1：CSV → PPTX ==========
    print("=" * 50)
    print("阶段 1: CSV → PPTX")
    print("=" * 50)

    df = pd.read_csv(args.csv, dtype={"id": str, "number": str})
    has_en_column = "en_name" in df.columns
    generated = []

    for _, row in df.iterrows():
        uid = clean_cell(row["id"])
        cn_name = str(row["cn_name"]).strip()
        num = clean_cell(row["number"])
        en_name = resolve_en_name(
            cn_name,
            row.get("en_name") if has_en_column else "",
            has_en_column
        )

        prs = fill_template(args.template, {
            "{{ID}}": uid,
            "{{CN_NAME}}": cn_name,
            "{{EN_NAME}}": en_name
        })

        fname = f"{num}_{uid}_{cn_name}.pptx"
        path = os.path.join(args.out_pptx, fname)
        prs.save(path)
        generated.append(path)
        print(f"  {fname}")

    print(f"\nPPTX 生成完成: {len(generated)} 个文件 → {args.out_pptx}\n")

    # ========== 阶段 2：PPTX → PNG（可选） ==========
    if args.no_png:
        print("已跳过 PNG 转换（--no-png）")
        print("全部完成！")
        return

    print("=" * 50)
    print("阶段 2: PPTX → PNG")
    print("=" * 50)

    # 检查依赖
    if not os.path.exists(SOFFICE_PATH):
        print("[跳过] 未安装 LibreOffice，无法转换 PNG")
        print(f"  macOS 安装: brew install --cask libreoffice")
        print("全部完成！")
        return

    pdftoppm = POPPLER_BIN / "pdftoppm"
    if not pdftoppm.exists():
        print("[跳过] 未安装 poppler，无法转换 PNG")
        print(f"  macOS 安装: brew install poppler")
        print("全部完成！")
        return

    os.makedirs(args.out_png, exist_ok=True)
    success = 0
    for pptx in sorted(Path(args.out_pptx).glob("*.pptx")):
        png_path = os.path.join(args.out_png, f"{pptx.stem}.png")
        print(f"  {pptx.name} → {os.path.basename(png_path)}", end="")
        if convert_to_png(pptx, png_path, pdftoppm):
            print("  ✓")
            success += 1
        else:
            print("  ✗")

    print(f"\nPNG 转换完成: {success}/{len(generated)} → {args.out_png}")
    print("全部完成！")


if __name__ == "__main__":
    main()
