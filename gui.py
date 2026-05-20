#!/usr/bin/env python3
"""
证书批量生成器 — 桌面版
双击运行或 python gui.py 启动可视化界面。
"""

import os
import subprocess
import sys
import tempfile
import threading
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pypinyin import lazy_pinyin, load_phrases_dict


# ====================== 工具函数 ======================

def is_chinese(text):
    if not text:
        return False
    return any('一' <= ch <= '鿿' for ch in str(text))


def load_surname_dict(csv_path):
    if not os.path.exists(csv_path):
        return {}
    df = pd.read_csv(csv_path, dtype=str)
    phrases = {}
    for _, row in df.iterrows():
        char = row["character"].strip()
        pinyin_str = row["surname_pinyin"].strip()
        phrases[char] = [[s] for s in pinyin_str.split()]
    return phrases


def chinese_to_pinyin(name):
    if not name or str(name).strip() == "":
        return ""
    pinyin_list = lazy_pinyin(str(name), strict=False)
    return " ".join(w.capitalize() for w in pinyin_list)


def fill_template(pptx_path, replacements):
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for key, value in replacements.items():
                            run.text = run.text.replace(key, value)
    return prs


def clean_cell(val):
    s = str(val).strip()
    return s[:-2] if s.endswith(".0") else s


# ====================== PNG 转换 ======================

SOFFICE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
POPPLER_BIN = Path("/opt/homebrew/opt/poppler/bin")


def check_png_deps():
    return os.path.exists(SOFFICE_PATH) and (POPPLER_BIN / "pdftoppm").exists()


def convert_to_png(pptx_path, output_png):
    pdftoppm = POPPLER_BIN / "pdftoppm"
    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_dir = Path(tmpdir)
        result = subprocess.run(
            [SOFFICE_PATH, "--headless", "--convert-to", "pdf",
             "--outdir", str(pdf_dir), str(pptx_path)],
            capture_output=True, text=True
        )
        if result.returncode != 0:
            return False
        pdf_files = list(pdf_dir.glob("*.pdf"))
        if not pdf_files:
            return False
        prefix = Path(tmpdir) / "page"
        result = subprocess.run(
            [str(pdftoppm), "-png", "-r", "200", "-f", "1", "-l", "1",
             str(pdf_files[0]), str(prefix)],
            capture_output=True, text=True
        )
        if result.returncode != 0:
            return False
        generated = list(Path(tmpdir).glob("page-*.png"))
        if not generated:
            return False
        generated[0].rename(output_png)
        return True


# ====================== GUI ======================

class App:
    def __init__(self, root):
        self.root = root
        self.root.title("证书批量生成器")
        self.root.geometry("640x500")
        self.root.resizable(False, False)

        # 样式
        style = ttk.Style()
        style.theme_use("clam")

        # ----- 标题 -----
        title = ttk.Label(root, text="证书批量生成器", font=("Helvetica", 18, "bold"))
        title.pack(pady=(20, 5))
        ttk.Label(root, text="上传 CSV + PPTX 模板，一键批量生成证书", foreground="#666").pack()

        # ----- 分隔线 -----
        ttk.Separator(root, orient="horizontal").pack(fill="x", padx=30, pady=15)

        # ----- CSV 文件选择 -----
        csv_frame = ttk.Frame(root)
        csv_frame.pack(fill="x", padx=40, pady=5)
        ttk.Label(csv_frame, text="CSV 数据文件", width=14).pack(side="left")
        self.csv_var = tk.StringVar()
        ttk.Entry(csv_frame, textvariable=self.csv_var, width=40).pack(side="left", padx=(5, 5))
        ttk.Button(csv_frame, text="选择...", command=self.select_csv, width=8).pack(side="left")

        # ----- PPTX 模板选择 -----
        pptx_frame = ttk.Frame(root)
        pptx_frame.pack(fill="x", padx=40, pady=5)
        ttk.Label(pptx_frame, text="PPTX 模板文件", width=14).pack(side="left")
        self.pptx_var = tk.StringVar()
        ttk.Entry(pptx_frame, textvariable=self.pptx_var, width=40).pack(side="left", padx=(5, 5))
        ttk.Button(pptx_frame, text="选择...", command=self.select_pptx, width=8).pack(side="left")

        # ----- 输出目录选择 -----
        out_frame = ttk.Frame(root)
        out_frame.pack(fill="x", padx=40, pady=5)
        ttk.Label(out_frame, text="输出到", width=14).pack(side="left")
        self.out_var = tk.StringVar(value=os.path.join(os.path.expanduser("~"), "Desktop", "Certificates"))
        ttk.Entry(out_frame, textvariable=self.out_var, width=40).pack(side="left", padx=(5, 5))
        ttk.Button(out_frame, text="选择...", command=self.select_output, width=8).pack(side="left")

        # ----- 选项 -----
        opt_frame = ttk.Frame(root)
        opt_frame.pack(fill="x", padx=40, pady=10)
        self.png_var = tk.BooleanVar(value=check_png_deps())
        png_state = "normal" if check_png_deps() else "disabled"
        ttk.Checkbutton(opt_frame, text="同时生成 PNG 图片", variable=self.png_var, state=png_state).pack(side="left")
        if not check_png_deps():
            ttk.Label(opt_frame, text="（需安装 LibreOffice）", foreground="#999").pack(side="left", padx=5)

        # ----- 进度条 -----
        self.progress = ttk.Progressbar(root, mode="determinate", length=560)
        self.progress.pack(pady=(10, 5))

        self.status_var = tk.StringVar(value="等待开始...")
        ttk.Label(root, textvariable=self.status_var, foreground="#666").pack()

        # ----- 按钮 -----
        btn_frame = ttk.Frame(root)
        btn_frame.pack(pady=20)
        self.run_btn = ttk.Button(btn_frame, text="开始生成", command=self.start, width=20)
        self.run_btn.pack(side="left", padx=5)
        ttk.Button(btn_frame, text="打开输出目录", command=self.open_output, width=15).pack(side="left", padx=5)

        # 加载姓氏字典
        script_dir = os.path.dirname(os.path.abspath(__file__))
        surname_csv = os.path.join(script_dir, "surname_pinyin.csv")
        if os.path.exists(surname_csv):
            load_phrases_dict(load_surname_dict(surname_csv))

    def select_csv(self):
        path = filedialog.askopenfilename(
            title="选择 CSV 数据文件",
            filetypes=[("CSV 文件", "*.csv"), ("所有文件", "*.*")]
        )
        if path:
            self.csv_var.set(path)

    def select_pptx(self):
        path = filedialog.askopenfilename(
            title="选择 PPTX 模板文件",
            filetypes=[("PPTX 文件", "*.pptx"), ("所有文件", "*.*")]
        )
        if path:
            self.pptx_var.set(path)

    def select_output(self):
        path = filedialog.askdirectory(title="选择输出目录")
        if path:
            self.out_var.set(path)

    def open_output(self):
        out = self.out_var.get()
        if os.path.exists(out):
            subprocess.run(["open", out])
        else:
            messagebox.showinfo("提示", "输出目录尚不存在，请先生成文件。")

    def start(self):
        csv_path = self.csv_var.get()
        pptx_path = self.pptx_var.get()
        out_dir = self.out_var.get()

        # 校验
        if not csv_path or not os.path.exists(csv_path):
            messagebox.showerror("错误", "请选择有效的 CSV 文件")
            return
        if not pptx_path or not os.path.exists(pptx_path):
            messagebox.showerror("错误", "请选择有效的 PPTX 模板文件")
            return

        self.run_btn.config(state="disabled")
        self.status_var.set("正在生成...")
        self.progress["value"] = 0

        thread = threading.Thread(target=self.run, args=(csv_path, pptx_path, out_dir))
        thread.daemon = True
        thread.start()

    def run(self, csv_path, pptx_path, out_dir):
        try:
            out_pptx = os.path.join(out_dir, "pptx")
            out_png = os.path.join(out_dir, "png")
            os.makedirs(out_pptx, exist_ok=True)

            df = pd.read_csv(csv_path, dtype={"id": str, "number": str})
            has_en_column = "en_name" in df.columns
            total = len(df)
            generated = []

            for i, (_, row) in enumerate(df.iterrows()):
                uid = clean_cell(row["id"])
                cn_name = str(row["cn_name"]).strip()
                num = clean_cell(row["number"])

                en_name = ""
                if is_chinese(cn_name):
                    if has_en_column:
                        val = row["en_name"]
                        if pd.notna(val) and str(val).strip() != "":
                            en_name = str(val).strip()
                    if en_name == "":
                        en_name = chinese_to_pinyin(cn_name)

                prs = fill_template(pptx_path, {
                    "{{ID}}": uid,
                    "{{CN_NAME}}": cn_name,
                    "{{EN_NAME}}": en_name
                })

                fname = f"{num}_{uid}_{cn_name}.pptx"
                path = os.path.join(out_pptx, fname)
                prs.save(path)
                generated.append(path)

                self.root.after(0, self.update_progress, i + 1, total, f"PPTX: {i + 1}/{total}")

            # PNG 转换
            do_png = self.png_var.get() and check_png_deps()
            if do_png:
                os.makedirs(out_png, exist_ok=True)
                success = 0
                for j, pptx in enumerate(generated):
                    png_path = os.path.join(out_png, f"{Path(pptx).stem}.png")
                    ok = convert_to_png(pptx, png_path)
                    if ok:
                        success += 1
                    self.root.after(0, self.update_progress, j + 1, total,
                                    f"PNG: {j + 1}/{total}")
                self.root.after(0, self.done, f"完成！PPTX {total} 个" +
                                (f"，PNG {success}/{total} 个 → {out_dir}" if do_png else f" → {out_dir}"))
            else:
                self.root.after(0, self.done, f"完成！PPTX {total} 个 → {out_dir}")

        except Exception as e:
            self.root.after(0, self.done, f"出错: {e}")

    def update_progress(self, done, total, msg):
        self.progress["value"] = (done / total) * 100
        self.status_var.set(msg)

    def done(self, msg):
        self.progress["value"] = 100
        self.status_var.set(msg)
        self.run_btn.config(state="normal")
        self.root.after(1000, lambda: messagebox.showinfo("完成", msg))


if __name__ == "__main__":
    root = tk.Tk()
    app = App(root)
    root.mainloop()
